import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from src import project_workspace
from src.business_query_history_service import (
    build_saved_business_query_report_context,
    get_saved_queries_for_dataset,
    save_query,
)
from src.business_analysis import request_management_summary
from src.exporter import (
    build_saved_business_queries_ai_context,
    export_executive_ppt,
    export_full_excel_report,
    export_ppt_from_template,
    export_word_from_template,
    generate_ai_periodic_report,
)


class SavedBusinessQueryReportTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.project_root = Path(self.temp_dir.name) / "projects"
        self.root_patch = patch.object(
            project_workspace,
            "PROJECT_ROOT",
            self.project_root,
        )
        self.root_patch.start()
        self.project = project_workspace.create_project("Saved Query Reports")
        self.project_id = self.project["project_id"]
        self.current_df = pd.DataFrame(
            {
                "成交日期": pd.to_datetime(["2026-01-01", "2026-02-01"]),
                "成交金额": [100.0, 200.0],
            }
        )

    def tearDown(self):
        self.root_patch.stop()
        self.temp_dir.cleanup()

    def _save_query(self, question, dataset_key, identifier):
        result = pd.DataFrame(
            {
                "排名": [1],
                "销售工号": [identifier],
                "订单数": [5],
            }
        )
        return save_query(
            self.project_id,
            question,
            {"dimension": "销售工号", "metric": "订单数", "limit": 1},
            result,
            f"{question}的解读",
            dataset_key,
            f"{dataset_key}.csv",
            is_identifier_dimension=True,
        )

    def _context_for(self, dataset_key):
        return build_saved_business_query_report_context(
            get_saved_queries_for_dataset(self.project_id, dataset_key)
        )

    def test_excel_multiple_queries_have_summary_and_individual_sheets(self):
        self._save_query("当前数据集问题一", "dataset-a", "100000003000000001")
        self._save_query("当前数据集问题二", "dataset-a", "100000003000000002")
        context = self._context_for("dataset-a")

        workbook_bytes = export_full_excel_report(
            self.current_df,
            saved_business_queries=context,
        )
        workbook = load_workbook(BytesIO(workbook_bytes), data_only=True)

        self.assertIn("Saved Business Queries", workbook.sheetnames)
        self.assertIn("Query_01", workbook.sheetnames)
        self.assertIn("Query_02", workbook.sheetnames)
        summary_values = [
            cell.value
            for row in workbook["Saved Business Queries"].iter_rows()
            for cell in row
        ]
        self.assertIn("当前数据集问题一", summary_values)
        self.assertIn("当前数据集问题二", summary_values)
        query_values = [
            cell.value
            for row in workbook["Query_01"].iter_rows()
            for cell in row
        ]
        self.assertIn("100000003000000001", query_values)
        workbook.close()

    def test_empty_saved_queries_do_not_break_excel_word_or_ppt(self):
        excel_bytes = export_full_excel_report(
            self.current_df,
            saved_business_queries=[],
        )
        word_bytes = export_word_from_template(
            self.current_df,
            template_file=None,
            output_path=None,
            context={"saved_business_queries": []},
        )
        ppt_bytes = export_ppt_from_template(
            self.current_df,
            template_file=None,
            output_path=None,
            context={"saved_business_queries": []},
        )

        workbook = load_workbook(BytesIO(excel_bytes), data_only=True)
        self.assertEqual(
            workbook["Saved Business Queries"]["A2"].value,
            "当前数据集暂无已保存业务查询。",
        )
        self.assertFalse(any(name.startswith("Query_") for name in workbook.sheetnames))
        workbook.close()
        document_text = "\n".join(
            paragraph.text for paragraph in Document(BytesIO(word_bytes)).paragraphs
        )
        self.assertIn("当前分析数据集暂无已保存业务查询结果。", document_text)
        presentation_text = _presentation_text(Presentation(BytesIO(ppt_bytes)))
        self.assertIn("当前数据集暂无已保存业务查询结果。", presentation_text)

    def test_word_ppt_and_ai_only_receive_current_dataset_queries(self):
        self._save_query("当前数据集问题", "dataset-a", "100000003000000001")
        self._save_query("其他数据集问题", "dataset-b", "200000003000000001")
        context = self._context_for("dataset-a")

        word_bytes = export_word_from_template(
            self.current_df,
            template_file=None,
            output_path=None,
            context={"saved_business_queries": context},
        )
        ppt_bytes = export_ppt_from_template(
            self.current_df,
            template_file=None,
            output_path=None,
            context={"saved_business_queries": context},
        )
        executive_bytes = export_executive_ppt(
            self.current_df,
            saved_business_queries=context,
        )
        prompts = []
        result = generate_ai_periodic_report(
            self.current_df,
            "月报",
            "成交日期",
            ["成交金额"],
            ai_client=lambda prompt: prompts.append(prompt) or "完成",
            saved_business_queries=context,
        )

        word_document = Document(BytesIO(word_bytes))
        word_text = "\n".join(
            [paragraph.text for paragraph in word_document.paragraphs]
            + [
                cell.text
                for table in word_document.tables
                for row in table.rows
                for cell in row.cells
            ]
        )
        ppt_text = _presentation_text(Presentation(BytesIO(ppt_bytes)))
        executive_text = _presentation_text(Presentation(BytesIO(executive_bytes)))

        for report_text in (word_text, ppt_text, executive_text, prompts[0]):
            self.assertIn("当前数据集问题", report_text)
            self.assertNotIn("其他数据集问题", report_text)
        self.assertIn("100000003000000001", word_text)
        self.assertEqual(result, "完成")

    def test_ai_context_caps_each_saved_result_at_twenty_rows(self):
        result = pd.DataFrame(
            {
                "排名": range(1, 26),
                "小组": [f"小组-{index}" for index in range(1, 26)],
                "成交金额": range(25, 0, -1),
            }
        )
        context = [
            {
                "question": "全部小组排名",
                "metric": "成交金额",
                "dimension": "小组",
                "result_df": result,
                "result_row_count": 25,
                "explanation": "测试解读",
            }
        ]

        ai_context = build_saved_business_queries_ai_context(context, max_rows=20)

        self.assertEqual(ai_context[0]["result_row_count"], 25)
        self.assertEqual(len(ai_context[0]["result_rows"]), 20)

    def test_ai_prompt_states_when_no_saved_queries_exist(self):
        prompts = []

        generate_ai_periodic_report(
            self.current_df,
            "月报",
            "成交日期",
            ["成交金额"],
            ai_client=lambda prompt: prompts.append(prompt) or "完成",
            saved_business_queries=[],
        )

        self.assertIn("当前数据集没有用户保存的业务查询结果。", prompts[0])

    def test_incomplete_saved_record_does_not_break_report_generators(self):
        incomplete_context = [
            {
                "question": "旧记录",
                "metric": "-",
                "dimension": "-",
                "result_rows": [],
                "result_columns": [],
                "explanation": "-",
                "dataset_name": "-",
            }
        ]

        excel_bytes = export_full_excel_report(
            self.current_df,
            saved_business_queries=incomplete_context,
        )
        word_bytes = export_word_from_template(
            self.current_df,
            output_path=None,
            context={"saved_business_queries": incomplete_context},
        )
        ppt_bytes = export_ppt_from_template(
            self.current_df,
            output_path=None,
            context={"saved_business_queries": incomplete_context},
        )

        self.assertTrue(excel_bytes)
        self.assertTrue(word_bytes)
        self.assertTrue(ppt_bytes)

    def test_management_summary_prompt_marks_saved_results_as_confirmed(self):
        payload = {
            "用户确认并保存的业务查询结果": [
                {"question": "当前数据集问题", "result_rows": [{"订单数": 5}]}
            ]
        }
        with (
            patch("src.business_analysis._request_completion", return_value={}) as request,
            patch("src.business_analysis._extract_text", return_value="完成"),
        ):
            result = request_management_summary(
                payload,
                "test-key",
                "test-model",
                "https://example.com/v1",
            )

        prompt = request.call_args.args[0]
        self.assertEqual(result, "完成")
        self.assertIn("用户主动保存的分析快照", prompt)
        self.assertIn("当前数据集问题", prompt)


def _presentation_text(presentation) -> str:
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


if __name__ == "__main__":
    unittest.main()
