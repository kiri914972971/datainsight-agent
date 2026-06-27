from __future__ import annotations

import json
import math
import posixpath
import re
from copy import copy
from datetime import date, datetime
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET
from zipfile import ZIP_DEFLATED, ZipFile

import pandas as pd
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows

from src.dashboard_exporter import create_excel_dashboard
from src.eda_ai_complete import _extract_text, _request_completion
from src.template_manager import get_last_sheet_name


REPORT_TEMPLATE_PLACEHOLDERS = [
    "{{report_title}}",
    "{{date_range}}",
    "{{data_overview}}",
    "{{quality_summary}}",
    "{{numeric_summary}}",
    "{{categorical_summary}}",
    "{{correlation_summary}}",
    "{{business_summary}}",
    "{{kpi_summary}}",
    "{{trend_summary}}",
    "{{topn_summary}}",
    "{{ai_insights}}",
    "{{risks}}",
    "{{recommendations}}",
]

EXCEL_DATA_SOURCE_SHEET_NAMES = [
    "Raw Data",
    "原始数据",
    "Data",
    "数据源",
    "销售数据",
    "明细数据",
    "Processed Data",
]

_MAIN_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CONTENT_TYPE_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_WORKSHEET_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet"
)
_WORKSHEET_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"
)

ET.register_namespace("", _MAIN_NS)
ET.register_namespace("r", _REL_NS)


def export_processed_data_excel(
    df: pd.DataFrame,
    original_df: pd.DataFrame | None = None,
    quality_summary: Any = None,
    repair_suggestions: Any = None,
) -> bytes:
    """Export the current processed dataset, optionally with supporting sheets."""
    sheets: list[tuple[str, Any]] = []
    if original_df is not None:
        sheets.append(("Original Data", original_df))
    sheets.append(("Processed Data", df))
    sheets.append(("Data Quality Summary", quality_summary))
    sheets.append(("Repair Suggestions", repair_suggestions))
    return _export_excel_sheets(sheets)


def export_full_excel_report(
    current_df: pd.DataFrame,
    quality_summary: Any = None,
    numeric_summary: Any = None,
    categorical_summary: Any = None,
    business_result: Any = None,
) -> bytes:
    """Export the full V1 Excel workbook, skipping unavailable results."""
    return _export_excel_sheets(
        [
            ("Processed Data", current_df),
            ("Data Quality Summary", quality_summary),
            ("Numeric Summary", numeric_summary),
            ("Categorical Summary", categorical_summary),
            ("Business Analysis Result", business_result),
        ]
    )


def export_excel_dashboard_from_template(
    current_df: pd.DataFrame,
    template_path: str | Path | bytes | None = None,
    output_path: str | Path | None = None,
    quality_summary: Any = None,
    business_summary: Any = None,
    kpi_summary: Any = None,
) -> bytes:
    """Backward-compatible entry point for the generated Excel Dashboard."""
    return create_excel_dashboard(current_df, output_path=output_path)


def _export_excel_dashboard_from_template_package_legacy(
    current_df: pd.DataFrame,
    template_path: str | Path | bytes,
    output_path: str | Path | None,
    quality_summary: Any = None,
    business_summary: Any = None,
    kpi_summary: Any = None,
) -> bytes:
    """Update the complete Dashboard workbook while preserving all template parts."""
    template_bytes = _read_excel_template_bytes(template_path)
    frame = _excel_safe_dataframe(current_df)
    replacements: dict[str, bytes] = {}
    additions: dict[str, bytes] = {}

    with ZipFile(BytesIO(template_bytes), "r") as template_package:
        sheet_paths = _xlsx_sheet_paths(template_package)
        data_source_sheet = _detect_xlsx_data_source_sheet(
            template_package,
            sheet_paths,
        )
        source_sheet_found = data_source_sheet is not None

        if data_source_sheet is None:
            data_source_sheet = "Processed Data"
            source_sheet_path = _add_xlsx_sheet(
                template_package,
                replacements,
                additions,
                data_source_sheet,
            )
            source_table_paths: list[str] = []
        else:
            source_sheet_path = sheet_paths[data_source_sheet]
            source_table_paths = _xlsx_sheet_table_paths(
                template_package,
                source_sheet_path,
            )

        source_sheet_xml = (
            replacements.get(source_sheet_path)
            or additions.get(source_sheet_path)
            or template_package.read(source_sheet_path)
        )
        replacements[source_sheet_path] = _replace_xlsx_sheet_data(
            source_sheet_xml,
            frame,
        )

        table_ref = _dataframe_excel_ref(frame)
        for table_path in source_table_paths:
            replacements[table_path] = _update_xlsx_table(
                template_package.read(table_path),
                table_ref,
                list(frame.columns),
            )

        _set_xlsx_refresh_on_open(template_package, replacements)
        _ensure_xlsx_readme_sheet(
            template_package,
            replacements,
            additions,
            data_source_sheet,
            source_sheet_found,
        )
        result = _write_updated_xlsx_package(
            template_package,
            replacements,
            additions,
        )

    return _finalize_binary_export(result, output_path)


def inspect_excel_dashboard_template(template_path: str | Path | bytes) -> dict[str, Any]:
    """Return template structure and detected data-source details for the UI."""
    template_bytes = _read_excel_template_bytes(template_path)
    with ZipFile(BytesIO(template_bytes), "r") as package:
        sheet_paths = _xlsx_sheet_paths(package)
        source_sheet = _detect_xlsx_data_source_sheet(package, sheet_paths)
        names = package.namelist()
        return {
            "sheet_names": list(sheet_paths),
            "source_sheet": source_sheet,
            "source_sheet_found": source_sheet is not None,
            "has_pivot_tables": any(name.startswith("xl/pivotTables/") for name in names),
            "has_slicers": any(
                name.startswith(("xl/slicers/", "xl/slicerCaches/"))
                for name in names
            ),
        }


def _export_excel_dashboard_from_template_static_legacy(
    current_df: pd.DataFrame,
    template_path: str | Path | bytes,
    output_path: str | Path | None,
    quality_summary: Any = None,
    business_summary: Any = None,
    kpi_summary: Any = None,
) -> bytes:
    """Create a static dashboard workbook from the final sheet of an Excel template."""
    if isinstance(template_path, bytes):
        template_source = BytesIO(template_path)
    else:
        template_path = Path(template_path)
        if not template_path.exists():
            raise FileNotFoundError(f"未找到 Excel Dashboard 模板：{template_path}")
        template_source = template_path

    template_workbook = load_workbook(
        template_source,
        data_only=False,
        read_only=False,
        keep_links=False,
    )
    dashboard_sheet_name = get_last_sheet_name(template_workbook)
    source_sheet = template_workbook[dashboard_sheet_name]

    output_workbook = Workbook()
    output_workbook.remove(output_workbook.active)
    dashboard_sheet = output_workbook.create_sheet(dashboard_sheet_name)
    _copy_dashboard_sheet(source_sheet, dashboard_sheet)

    placeholder_values = _dashboard_placeholder_values(
        current_df,
        quality_summary,
        business_summary,
        kpi_summary,
    )
    replaced_placeholders = _replace_dashboard_placeholders(
        dashboard_sheet,
        placeholder_values,
    )
    _write_dashboard_fallback_summary(
        dashboard_sheet,
        placeholder_values,
        replaced_placeholders,
    )

    _append_dataframe_sheet(output_workbook, "Processed Data", current_df)
    _append_content_sheet(
        output_workbook,
        "Data Quality Summary",
        quality_summary,
        "暂无数据质量结果",
    )
    _append_dataframe_sheet(
        output_workbook,
        "Exploration Summary",
        _build_exploration_summary(current_df),
    )
    _append_content_sheet(
        output_workbook,
        "Business Analysis Summary",
        business_summary,
        "暂无业务分析结果",
    )

    buffer = BytesIO()
    output_workbook.save(buffer)
    template_workbook.close()
    output_workbook.close()
    result = buffer.getvalue()

    if output_path is not None:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(result)
    return result


def export_word_report(
    current_df: pd.DataFrame,
    quality_summary: Any = None,
    numeric_summary: Any = None,
    categorical_summary: Any = None,
    business_summary: Any = None,
    ai_summary: str | None = None,
) -> bytes:
    document = Document()
    document.add_heading("DataInsight Agent 数据分析报告", 0)

    document.add_heading("1. 数据概况", level=1)
    memory_mb = current_df.memory_usage(deep=True).sum() / 1024 / 1024
    for line in (
        f"行数：{len(current_df):,}",
        f"列数：{len(current_df.columns):,}",
        f"重复行：{int(current_df.duplicated().sum()):,}",
        f"内存占用：{memory_mb:.2f} MB",
    ):
        document.add_paragraph(line, style="List Bullet")

    document.add_heading("2. 数据质量摘要", level=1)
    _add_content(document, quality_summary)

    document.add_heading("3. 探索性分析摘要", level=1)
    document.add_heading("数值字段统计", level=2)
    _add_content(document, numeric_summary)
    document.add_heading("类别字段统计", level=2)
    _add_content(document, categorical_summary)
    document.add_paragraph("相关性分析摘要请结合应用中的相关分析热力图与高相关字段对阅读。")

    document.add_heading("4. 业务分析摘要", level=1)
    _add_content(document, business_summary)

    document.add_heading("5. AI总结", level=1)
    document.add_paragraph(ai_summary or "尚未生成 AI 总结。")

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def export_word_from_template(
    current_df: pd.DataFrame,
    template_file=None,
    output_path: str | Path | None = "data_insight_report.docx",
    context: dict | None = None,
) -> bytes:
    """Generate a Word report by replacing stable placeholders in a template."""
    replacements = _build_report_template_context(current_df, context)
    if template_file is None:
        document = _build_default_word_report(replacements)
    else:
        document = Document(_binary_template_source(template_file))
        _replace_word_document_placeholders(document, replacements)

    buffer = BytesIO()
    document.save(buffer)
    return _finalize_binary_export(buffer.getvalue(), output_path)


def export_ppt_report(
    current_df: pd.DataFrame,
    quality_summary: Any = None,
    business_summary: Any = None,
    ai_summary: str | None = None,
) -> bytes:
    """Generate the V2 text-first seven-slide presentation."""
    Presentation, Inches, Pt = _pptx_dependencies()
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    _add_ppt_title_slide(presentation, "DataInsight Agent 数据分析报告", "基于当前处理后数据生成")
    _add_ppt_bullets(
        presentation,
        "数据概况",
        [
            f"数据规模：{len(current_df):,} 行 × {len(current_df.columns):,} 列",
            f"重复行：{int(current_df.duplicated().sum()):,}",
            f"缺失值总数：{int(current_df.isna().sum().sum()):,}",
        ],
        Pt,
    )
    quality_lines = _content_lines(quality_summary, limit=7) or ["暂无数据质量摘要。"]
    _add_ppt_bullets(presentation, "数据质量", quality_lines, Pt)

    business_lines = _content_lines(business_summary, limit=10)
    _add_ppt_bullets(presentation, "核心 KPI", business_lines[:5] or ["暂无核心 KPI 结果。"], Pt)
    _add_ppt_bullets(presentation, "时间趋势", business_lines[5:10] or ["请在业务分析中生成时间趋势。"], Pt)
    _add_ppt_bullets(presentation, "维度对比 / Top N", business_lines[:8] or ["请在业务分析中生成维度排行。"], Pt)
    _add_ppt_bullets(
        presentation,
        "AI 总结与建议",
        _text_lines(ai_summary, limit=10) or ["尚未生成 AI 总结。"],
        Pt,
    )
    return _save_presentation(presentation)


def export_ppt_from_template(
    current_df: pd.DataFrame,
    template_file=None,
    output_path: str | Path | None = "data_insight_presentation.pptx",
    context: dict | None = None,
) -> bytes:
    """Generate a PowerPoint report by replacing placeholders in slide text."""
    replacements = _build_report_template_context(current_df, context)
    Presentation, Inches, Pt = _pptx_dependencies()
    if template_file is None:
        presentation = _build_default_ppt_report(replacements, Presentation, Inches, Pt)
    else:
        presentation = Presentation(_binary_template_source(template_file))
        _replace_ppt_placeholders(presentation, replacements)

    result = _save_presentation(presentation)
    return _finalize_binary_export(result, output_path)


def generate_ai_periodic_report(
    current_df: pd.DataFrame,
    period_type: str,
    date_column: str,
    metric_columns: list[str],
    ai_client=None,
    *,
    api_key: str | None = None,
    model: str = "gpt-5.4-mini",
    base_url: str = "https://api.openai.com/v1",
) -> str:
    if date_column not in current_df.columns:
        raise KeyError(f"日期字段不存在：{date_column}")
    valid_metrics = [column for column in metric_columns if column in current_df.columns]
    if not valid_metrics:
        raise ValueError("请至少选择一个有效指标字段。")

    frequency = {"周报": "W", "月报": "M", "季度报告": "Q", "年度报告": "Y"}[period_type]
    frame = current_df[[date_column, *valid_metrics]].copy()
    frame[date_column] = pd.to_datetime(frame[date_column], errors="coerce")
    for column in valid_metrics:
        frame[column] = pd.to_numeric(frame[column], errors="coerce")
    frame = frame.dropna(subset=[date_column])
    frame["周期"] = frame[date_column].dt.to_period(frequency).astype(str)
    periodic = frame.groupby("周期", as_index=False)[valid_metrics].sum(min_count=1).tail(12)
    if periodic.empty:
        raise ValueError("当前日期字段没有可用于生成周期报告的有效数据。")

    prompt = f"""
你是一名经营分析顾问。根据下面的周期指标数据，用简体中文生成{period_type}。

严格使用以下标题：
1. 本期概况
2. 核心指标变化
3. 主要增长点
4. 主要风险
5. 下期建议

要求：
- 聚焦业务变化、风险和行动建议，不讨论缺失值、偏度、峰度或 IQR。
- 不虚构数据中不存在的原因；信息不足时明确提出需要验证的问题。
- 内容简洁，必须完整结束。

周期指标数据：
{json.dumps(periodic.to_dict("records"), ensure_ascii=False, default=str)}
""".strip()
    if ai_client is not None:
        return ai_client(prompt)
    if not api_key:
        raise ValueError("请先在侧边栏完成 AI 接入。")
    return _extract_text(_request_completion(prompt, api_key, model, base_url, timeout=90))


def export_executive_ppt(
    current_df: pd.DataFrame,
    kpi_summary: Any = None,
    ai_summary: str | None = None,
) -> bytes:
    """Generate the V4 five-slide management presentation."""
    Presentation, Inches, Pt = _pptx_dependencies()
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    ai_lines = _text_lines(ai_summary, limit=12)
    kpi_lines = _content_lines(kpi_summary, limit=8)

    _add_ppt_title_slide(presentation, "DataInsight Agent 管理层汇报", "决策摘要")
    _add_ppt_bullets(presentation, "管理层摘要", ai_lines[:5] or ["尚未生成管理层摘要。"], Pt)
    _add_ppt_bullets(presentation, "核心 KPI", kpi_lines or ["暂无核心 KPI 结果。"], Pt)
    _add_ppt_bullets(presentation, "增长亮点与主要风险", ai_lines[5:9] or ["请先生成 AI 业务总结。"], Pt)
    _add_ppt_bullets(presentation, "行动建议", ai_lines[9:12] or ["建议结合业务分析结果制定行动计划。"], Pt)
    return _save_presentation(presentation)


def _export_excel_sheets(sheets: list[tuple[str, Any]]) -> bytes:
    buffer = BytesIO()
    written = 0
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        for name, content in sheets:
            frame = _as_dataframe(content)
            if frame is None or frame.empty:
                continue
            _excel_safe_dataframe(frame).to_excel(writer, sheet_name=name[:31], index=False)
            written += 1
        if written == 0:
            pd.DataFrame({"说明": ["当前没有可导出的数据。"]}).to_excel(
                writer, sheet_name="Summary", index=False
            )
    return buffer.getvalue()


def _read_excel_template_bytes(template_path: str | Path | bytes) -> bytes:
    if isinstance(template_path, bytes):
        return template_path
    if hasattr(template_path, "getvalue"):
        return template_path.getvalue()
    if hasattr(template_path, "read"):
        return template_path.read()
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"未找到 Excel Dashboard 模板：{path}")
    return path.read_bytes()


def _xml_bytes(root: ET.Element) -> bytes:
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _package_part(
    package: ZipFile,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
    name: str,
) -> bytes:
    return replacements.get(name) or additions.get(name) or package.read(name)


def _resolve_xlsx_target(base_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(base_part), target))


def _xlsx_sheet_paths(package: ZipFile) -> dict[str, str]:
    workbook_root = ET.fromstring(package.read("xl/workbook.xml"))
    rels_root = ET.fromstring(package.read("xl/_rels/workbook.xml.rels"))
    relationship_paths = {
        relationship.get("Id"): _resolve_xlsx_target(
            "xl/workbook.xml",
            relationship.get("Target", ""),
        )
        for relationship in rels_root
        if relationship.get("Type") == _WORKSHEET_REL_TYPE
    }
    return {
        sheet.get("name", ""): relationship_paths[sheet.get(f"{{{_REL_NS}}}id")]
        for sheet in workbook_root.findall(f".//{{{_MAIN_NS}}}sheet")
        if sheet.get(f"{{{_REL_NS}}}id") in relationship_paths
    }


def _xlsx_sheet_table_paths(package: ZipFile, sheet_path: str) -> list[str]:
    rels_path = posixpath.join(
        posixpath.dirname(sheet_path),
        "_rels",
        f"{posixpath.basename(sheet_path)}.rels",
    )
    if rels_path not in package.namelist():
        return []
    rels_root = ET.fromstring(package.read(rels_path))
    return [
        _resolve_xlsx_target(sheet_path, relationship.get("Target", ""))
        for relationship in rels_root
        if relationship.get("Type", "").endswith("/table")
    ]


def _detect_xlsx_data_source_sheet(
    package: ZipFile,
    sheet_paths: dict[str, str],
) -> str | None:
    normalized_names = {name.strip().casefold(): name for name in sheet_paths}
    for candidate in EXCEL_DATA_SOURCE_SHEET_NAMES:
        if candidate.casefold() in normalized_names:
            return normalized_names[candidate.casefold()]

    table_to_sheet: dict[str, str] = {}
    for sheet_name, sheet_path in sheet_paths.items():
        for table_path in _xlsx_sheet_table_paths(package, sheet_path):
            table_root = ET.fromstring(package.read(table_path))
            for attribute in ("name", "displayName"):
                table_name = table_root.get(attribute)
                if table_name:
                    table_to_sheet[table_name.casefold()] = sheet_name

    for part_name in package.namelist():
        if not part_name.startswith("xl/pivotCache/pivotCacheDefinition"):
            continue
        cache_root = ET.fromstring(package.read(part_name))
        source = cache_root.find(f".//{{{_MAIN_NS}}}worksheetSource")
        if source is None:
            continue
        source_sheet = source.get("sheet")
        if source_sheet in sheet_paths:
            return source_sheet
        source_name = source.get("name", "").casefold()
        if source_name in table_to_sheet:
            return table_to_sheet[source_name]

    fallback_tokens = ("源数据", "原始", "明细", "raw data", "rawdata", "processed data")
    for sheet_name in sheet_paths:
        normalized = sheet_name.strip().casefold()
        if any(token in normalized for token in fallback_tokens):
            return sheet_name
    return None


def _dataframe_excel_ref(dataframe: pd.DataFrame) -> str:
    last_column = get_column_letter(max(1, len(dataframe.columns)))
    last_row = max(1, len(dataframe) + 1)
    return f"A1:{last_column}{last_row}"


def _replace_xlsx_sheet_data(sheet_xml: bytes, dataframe: pd.DataFrame) -> bytes:
    root = ET.fromstring(sheet_xml)
    sheet_data = root.find(f"{{{_MAIN_NS}}}sheetData")
    if sheet_data is None:
        sheet_data = ET.Element(f"{{{_MAIN_NS}}}sheetData")
        root.append(sheet_data)
    else:
        sheet_data.clear()

    rows = [list(dataframe.columns), *dataframe.itertuples(index=False, name=None)]
    for row_number, values in enumerate(rows, start=1):
        row = ET.SubElement(sheet_data, f"{{{_MAIN_NS}}}row", {"r": str(row_number)})
        for column_number, value in enumerate(values, start=1):
            _append_xlsx_cell(row, row_number, column_number, value)

    dimension = root.find(f"{{{_MAIN_NS}}}dimension")
    if dimension is None:
        dimension = ET.Element(f"{{{_MAIN_NS}}}dimension")
        root.insert(0, dimension)
    dimension.set("ref", _dataframe_excel_ref(dataframe))

    auto_filter = root.find(f"{{{_MAIN_NS}}}autoFilter")
    if auto_filter is not None:
        auto_filter.set("ref", _dataframe_excel_ref(dataframe))
    return _xml_bytes(root)


def _append_xlsx_cell(
    row: ET.Element,
    row_number: int,
    column_number: int,
    value: Any,
) -> None:
    coordinate = f"{get_column_letter(column_number)}{row_number}"
    if hasattr(value, "item") and not isinstance(value, (str, bytes)):
        try:
            value = value.item()
        except ValueError:
            pass
    try:
        if pd.isna(value):
            return
    except (TypeError, ValueError):
        pass

    cell = ET.SubElement(row, f"{{{_MAIN_NS}}}c", {"r": coordinate})
    if isinstance(value, bool):
        cell.set("t", "b")
        ET.SubElement(cell, f"{{{_MAIN_NS}}}v").text = "1" if value else "0"
    elif isinstance(value, (datetime, date)):
        cell.set("t", "d")
        ET.SubElement(cell, f"{{{_MAIN_NS}}}v").text = value.isoformat()
    elif isinstance(value, (int, float)) and not isinstance(value, bool):
        if isinstance(value, float) and not math.isfinite(value):
            _append_xlsx_inline_string(cell, str(value))
        else:
            ET.SubElement(cell, f"{{{_MAIN_NS}}}v").text = str(value)
    else:
        _append_xlsx_inline_string(cell, value)


def _append_xlsx_inline_string(cell: ET.Element, value: Any) -> None:
    cell.set("t", "inlineStr")
    inline = ET.SubElement(cell, f"{{{_MAIN_NS}}}is")
    text = ET.SubElement(inline, f"{{{_MAIN_NS}}}t")
    cleaned = ILLEGAL_CHARACTERS_RE.sub("", str(value))
    if cleaned != cleaned.strip():
        text.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    text.text = cleaned


def _update_xlsx_table(
    table_xml: bytes,
    table_ref: str,
    columns: list[Any],
) -> bytes:
    root = ET.fromstring(table_xml)
    root.set("ref", table_ref)
    auto_filter = root.find(f"{{{_MAIN_NS}}}autoFilter")
    if auto_filter is not None:
        auto_filter.set("ref", table_ref)
    table_columns = root.find(f"{{{_MAIN_NS}}}tableColumns")
    if table_columns is not None:
        table_columns.clear()
        table_columns.set("count", str(len(columns)))
        for index, column in enumerate(columns, start=1):
            ET.SubElement(
                table_columns,
                f"{{{_MAIN_NS}}}tableColumn",
                {"id": str(index), "name": str(column)},
            )
    return _xml_bytes(root)


def _new_xlsx_sheet_xml() -> bytes:
    root = ET.Element(f"{{{_MAIN_NS}}}worksheet")
    sheet_views = ET.SubElement(root, f"{{{_MAIN_NS}}}sheetViews")
    ET.SubElement(sheet_views, f"{{{_MAIN_NS}}}sheetView", {"workbookViewId": "0"})
    ET.SubElement(root, f"{{{_MAIN_NS}}}sheetFormatPr", {"defaultRowHeight": "15"})
    ET.SubElement(root, f"{{{_MAIN_NS}}}sheetData")
    return _xml_bytes(root)


def _next_relationship_id(rels_root: ET.Element) -> str:
    relationship_numbers = []
    for relationship in rels_root:
        match = re.fullmatch(r"rId(\d+)", relationship.get("Id", ""))
        if match:
            relationship_numbers.append(int(match.group(1)))
    return f"rId{max(relationship_numbers, default=0) + 1}"


def _next_worksheet_path(package: ZipFile, additions: dict[str, bytes]) -> str:
    worksheet_numbers = []
    for name in [*package.namelist(), *additions]:
        match = re.fullmatch(r"xl/worksheets/sheet(\d+)\.xml", name)
        if match:
            worksheet_numbers.append(int(match.group(1)))
    return f"xl/worksheets/sheet{max(worksheet_numbers, default=0) + 1}.xml"


def _add_xlsx_sheet(
    package: ZipFile,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
    sheet_name: str,
) -> str:
    workbook_path = "xl/workbook.xml"
    rels_path = "xl/_rels/workbook.xml.rels"
    content_types_path = "[Content_Types].xml"
    workbook_root = ET.fromstring(
        _package_part(package, replacements, additions, workbook_path)
    )
    rels_root = ET.fromstring(_package_part(package, replacements, additions, rels_path))
    content_types_root = ET.fromstring(
        _package_part(package, replacements, additions, content_types_path)
    )

    sheet_path = _next_worksheet_path(package, additions)
    relationship_id = _next_relationship_id(rels_root)
    sheets = workbook_root.find(f"{{{_MAIN_NS}}}sheets")
    sheet_ids = [int(sheet.get("sheetId", "0")) for sheet in sheets]
    ET.SubElement(
        sheets,
        f"{{{_MAIN_NS}}}sheet",
        {
            "name": sheet_name,
            "sheetId": str(max(sheet_ids, default=0) + 1),
            f"{{{_REL_NS}}}id": relationship_id,
        },
    )
    ET.SubElement(
        rels_root,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": relationship_id,
            "Type": _WORKSHEET_REL_TYPE,
            "Target": sheet_path.removeprefix("xl/"),
        },
    )
    ET.SubElement(
        content_types_root,
        f"{{{_CONTENT_TYPE_NS}}}Override",
        {
            "PartName": f"/{sheet_path}",
            "ContentType": _WORKSHEET_CONTENT_TYPE,
        },
    )
    replacements[workbook_path] = _xml_bytes(workbook_root)
    replacements[rels_path] = _xml_bytes(rels_root)
    replacements[content_types_path] = _xml_bytes(content_types_root)
    additions[sheet_path] = _new_xlsx_sheet_xml()
    return sheet_path


def _ensure_xlsx_readme_sheet(
    package: ZipFile,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
    data_source_sheet: str,
    source_sheet_found: bool,
) -> None:
    sheet_paths = _xlsx_sheet_paths_from_parts(package, replacements, additions)
    readme_path = sheet_paths.get("README")
    if readme_path is None:
        readme_path = _add_xlsx_sheet(package, replacements, additions, "README")
    readme_frame = pd.DataFrame(
        {
            "说明": [
                "本文件基于用户上传的 Excel Dashboard 模板生成。",
                f"当前数据已写入工作表：{data_source_sheet}",
                "如果 Dashboard 使用数据透视表或切片器，请打开 Excel 后点击：数据 -> 全部刷新。",
                "刷新操作将更新透视表、透视图和切片器。",
                "如果刷新后图表未更新，请检查透视表数据源是否指向当前数据源工作表或其中的数据表。",
                (
                    "系统未识别到模板原有数据源，已新增 Processed Data。请确认透视表数据源指向该工作表。"
                    if not source_sheet_found
                    else "系统已识别并更新模板原有数据源工作表。"
                ),
            ]
        }
    )
    base_xml = (
        replacements.get(readme_path)
        or additions.get(readme_path)
        or package.read(readme_path)
    )
    replacements[readme_path] = _replace_xlsx_sheet_data(base_xml, readme_frame)
    additions.pop(readme_path, None)


def _xlsx_sheet_paths_from_parts(
    package: ZipFile,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
) -> dict[str, str]:
    workbook_root = ET.fromstring(
        _package_part(package, replacements, additions, "xl/workbook.xml")
    )
    rels_root = ET.fromstring(
        _package_part(package, replacements, additions, "xl/_rels/workbook.xml.rels")
    )
    relationship_paths = {
        relationship.get("Id"): _resolve_xlsx_target(
            "xl/workbook.xml",
            relationship.get("Target", ""),
        )
        for relationship in rels_root
        if relationship.get("Type") == _WORKSHEET_REL_TYPE
    }
    return {
        sheet.get("name", ""): relationship_paths[sheet.get(f"{{{_REL_NS}}}id")]
        for sheet in workbook_root.findall(f".//{{{_MAIN_NS}}}sheet")
        if sheet.get(f"{{{_REL_NS}}}id") in relationship_paths
    }


def _set_xlsx_refresh_on_open(
    package: ZipFile,
    replacements: dict[str, bytes],
) -> None:
    workbook_path = "xl/workbook.xml"
    workbook_root = ET.fromstring(replacements.get(workbook_path) or package.read(workbook_path))
    calc_properties = workbook_root.find(f"{{{_MAIN_NS}}}calcPr")
    if calc_properties is None:
        calc_properties = ET.SubElement(workbook_root, f"{{{_MAIN_NS}}}calcPr")
    calc_properties.set("calcMode", "auto")
    calc_properties.set("fullCalcOnLoad", "1")
    calc_properties.set("forceFullCalc", "1")
    replacements[workbook_path] = _xml_bytes(workbook_root)

    for part_name in package.namelist():
        if not part_name.startswith("xl/pivotCache/pivotCacheDefinition"):
            continue
        cache_root = ET.fromstring(package.read(part_name))
        cache_root.set("refreshOnLoad", "1")
        cache_root.set("enableRefresh", "1")
        replacements[part_name] = _xml_bytes(cache_root)


def _write_updated_xlsx_package(
    package: ZipFile,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
) -> bytes:
    output = BytesIO()
    original_names = set(package.namelist())
    with ZipFile(output, "w", ZIP_DEFLATED) as result_package:
        for item in package.infolist():
            result_package.writestr(item, replacements.get(item.filename, package.read(item.filename)))
        for name in (set(additions) | set(replacements)) - original_names:
            result_package.writestr(
                name,
                replacements[name] if name in replacements else additions[name],
            )
    return output.getvalue()


def _copy_dashboard_sheet(source_sheet, target_sheet) -> None:
    target_sheet.sheet_format = copy(source_sheet.sheet_format)
    target_sheet.sheet_properties = copy(source_sheet.sheet_properties)
    target_sheet.page_margins = copy(source_sheet.page_margins)
    target_sheet.page_setup = copy(source_sheet.page_setup)
    target_sheet.print_options = copy(source_sheet.print_options)
    target_sheet.freeze_panes = source_sheet.freeze_panes
    target_sheet.sheet_view.showGridLines = source_sheet.sheet_view.showGridLines
    target_sheet.sheet_view.zoomScale = source_sheet.sheet_view.zoomScale

    for index, dimension in source_sheet.row_dimensions.items():
        target_dimension = target_sheet.row_dimensions[index]
        target_dimension.height = dimension.height
        target_dimension.hidden = dimension.hidden
        target_dimension.outlineLevel = dimension.outlineLevel
        target_dimension.collapsed = dimension.collapsed
        target_dimension.thickTop = dimension.thickTop
        target_dimension.thickBot = dimension.thickBot

    for key, dimension in source_sheet.column_dimensions.items():
        target_dimension = target_sheet.column_dimensions[key]
        target_dimension.width = dimension.width
        target_dimension.hidden = dimension.hidden
        target_dimension.bestFit = dimension.bestFit
        target_dimension.outlineLevel = dimension.outlineLevel
        target_dimension.collapsed = dimension.collapsed

    for source_row in source_sheet.iter_rows():
        for source_cell in source_row:
            target_cell = target_sheet[source_cell.coordinate]
            value = source_cell.value
            if isinstance(value, str) and value.startswith("="):
                value = None
            target_cell.value = value
            if source_cell.has_style:
                target_cell.font = copy(source_cell.font)
                target_cell.fill = copy(source_cell.fill)
                target_cell.border = copy(source_cell.border)
                target_cell.alignment = copy(source_cell.alignment)
                target_cell.number_format = source_cell.number_format
                target_cell.protection = copy(source_cell.protection)
            if source_cell.hyperlink:
                target_cell._hyperlink = copy(source_cell.hyperlink)
            if source_cell.comment:
                target_cell.comment = copy(source_cell.comment)

    for merged_range in source_sheet.merged_cells.ranges:
        target_sheet.merge_cells(str(merged_range))

    target_sheet.print_area = source_sheet.print_area
    target_sheet.print_title_cols = source_sheet.print_title_cols
    target_sheet.print_title_rows = source_sheet.print_title_rows


def _binary_template_source(template_file):
    if isinstance(template_file, bytes):
        return BytesIO(template_file)
    if isinstance(template_file, (str, Path)):
        path = Path(template_file)
        if not path.exists():
            raise FileNotFoundError(f"未找到报告模板：{path}")
        return path
    if hasattr(template_file, "getvalue"):
        return BytesIO(template_file.getvalue())
    if hasattr(template_file, "read"):
        return BytesIO(template_file.read())
    raise TypeError("无法读取报告模板。")


def _finalize_binary_export(result: bytes, output_path: str | Path | None) -> bytes:
    if output_path is not None:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(result)
    return result


def _build_report_template_context(
    current_df: pd.DataFrame,
    context: dict | None,
) -> dict[str, str]:
    memory_mb = current_df.memory_usage(deep=True).sum() / 1024 / 1024
    defaults = {
        "report_title": "DataInsight Agent 数据分析报告",
        "date_range": _detect_date_range(current_df),
        "data_overview": (
            f"行数：{len(current_df):,}\n"
            f"列数：{len(current_df.columns):,}\n"
            f"重复行：{int(current_df.duplicated().sum()):,}\n"
            f"缺失值总数：{int(current_df.isna().sum().sum()):,}\n"
            f"内存占用：{memory_mb:.2f} MB"
        ),
        "quality_summary": "暂无数据质量摘要。",
        "numeric_summary": "暂无数值字段统计。",
        "categorical_summary": "暂无类别字段统计。",
        "correlation_summary": "暂无相关性分析结果。",
        "business_summary": "暂无业务分析结果。",
        "kpi_summary": (
            f"数据行数：{len(current_df):,}\n"
            f"数据列数：{len(current_df.columns):,}"
        ),
        "trend_summary": "暂无趋势分析结果。",
        "topn_summary": "暂无 Top N 分析结果。",
        "ai_insights": "尚未生成 AI 洞察。",
        "risks": "暂无已确认风险，请结合业务场景判断。",
        "recommendations": "建议结合数据质量、探索性分析和业务分析结果制定行动计划。",
    }
    for key, value in (context or {}).items():
        normalized_key = str(key).strip("{}")
        if normalized_key in defaults and value is not None:
            defaults[normalized_key] = _format_template_content(value)
    return {f"{{{{{key}}}}}": str(value) for key, value in defaults.items()}


def _format_template_content(content: Any, max_lines: int = 30) -> str:
    if content is None:
        return ""
    if isinstance(content, str):
        return content
    frame = _as_dataframe(content)
    if frame is None or frame.empty:
        return ""
    lines = []
    for row in frame.head(max_lines).fillna("").to_dict("records"):
        lines.append("；".join(f"{key}：{value}" for key, value in row.items()))
    if len(frame) > max_lines:
        lines.append(f"仅展示前 {max_lines} 项，共 {len(frame)} 项。")
    return "\n".join(lines)


def _replace_run_placeholders(paragraph, replacements: dict[str, str]) -> int:
    runs = paragraph.runs
    if not runs:
        return 0

    replaced = 0
    for run in runs:
        for placeholder, replacement in replacements.items():
            count = run.text.count(placeholder)
            if count:
                run.text = run.text.replace(placeholder, replacement)
                replaced += count

    combined = "".join(run.text for run in runs)
    if not any(placeholder in combined for placeholder in replacements):
        return replaced

    updated = combined
    for placeholder, replacement in replacements.items():
        count = updated.count(placeholder)
        if count:
            updated = updated.replace(placeholder, replacement)
            replaced += count
    runs[0].text = updated
    for run in runs[1:]:
        run.text = ""
    return replaced


def _replace_word_document_placeholders(
    document: Document,
    replacements: dict[str, str],
) -> int:
    replaced = 0
    for paragraph in document.paragraphs:
        replaced += _replace_run_placeholders(paragraph, replacements)
    for table in document.tables:
        replaced += _replace_word_table_placeholders(table, replacements)
    for section in document.sections:
        for container in (section.header, section.footer):
            for paragraph in container.paragraphs:
                replaced += _replace_run_placeholders(paragraph, replacements)
            for table in container.tables:
                replaced += _replace_word_table_placeholders(table, replacements)
    return replaced


def _replace_word_table_placeholders(table, replacements: dict[str, str]) -> int:
    replaced = 0
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                replaced += _replace_run_placeholders(paragraph, replacements)
            for nested_table in cell.tables:
                replaced += _replace_word_table_placeholders(nested_table, replacements)
    return replaced


def _build_default_word_report(replacements: dict[str, str]) -> Document:
    document = Document()
    document.add_heading(replacements["{{report_title}}"], 0)
    document.add_paragraph(f"数据范围：{replacements['{{date_range}}']}")
    sections = [
        ("1. 数据概况", replacements["{{data_overview}}"]),
        ("2. 数据质量摘要", replacements["{{quality_summary}}"]),
        (
            "3. 探索性分析摘要",
            "\n\n".join(
                [
                    replacements["{{numeric_summary}}"],
                    replacements["{{categorical_summary}}"],
                    replacements["{{correlation_summary}}"],
                ]
            ),
        ),
        (
            "4. 业务分析摘要",
            "\n\n".join(
                [
                    replacements["{{kpi_summary}}"],
                    replacements["{{trend_summary}}"],
                    replacements["{{topn_summary}}"],
                    replacements["{{business_summary}}"],
                ]
            ),
        ),
        (
            "5. 风险与建议",
            "\n\n".join(
                [
                    f"风险\n{replacements['{{risks}}']}",
                    f"建议\n{replacements['{{recommendations}}']}",
                ]
            ),
        ),
    ]
    for heading, content in sections:
        document.add_heading(heading, level=1)
        document.add_paragraph(content or "暂无可用结果。")
    return document


def _replace_ppt_placeholders(presentation, replacements: dict[str, str]) -> int:
    replaced = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            replaced += _replace_ppt_shape_placeholders(shape, replacements)
    return replaced


def _replace_ppt_shape_placeholders(shape, replacements: dict[str, str]) -> int:
    replaced = 0
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            replaced += _replace_run_placeholders(paragraph, replacements)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replaced += _replace_run_placeholders(paragraph, replacements)
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            replaced += _replace_ppt_shape_placeholders(child_shape, replacements)
    return replaced


def _build_default_ppt_report(replacements, Presentation, Inches, Pt):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _add_ppt_title_slide(
        presentation,
        replacements["{{report_title}}"],
        f"数据范围：{replacements['{{date_range}}']}",
    )
    slides = [
        ("数据概况", replacements["{{data_overview}}"]),
        ("数据质量", replacements["{{quality_summary}}"]),
        ("核心 KPI", replacements["{{kpi_summary}}"]),
        (
            "业务分析",
            "\n".join(
                [
                    replacements["{{business_summary}}"],
                    replacements["{{trend_summary}}"],
                    replacements["{{topn_summary}}"],
                ]
            ),
        ),
        (
            "风险与建议",
            "\n".join(
                [
                    f"风险：{replacements['{{risks}}']}",
                    f"建议：{replacements['{{recommendations}}']}",
                ]
            ),
        ),
    ]
    for title, content in slides:
        _add_ppt_bullets(
            presentation,
            title,
            _text_lines(content, limit=10) or ["暂无可用结果。"],
            Pt,
        )
    return presentation


def _dashboard_placeholder_values(
    current_df: pd.DataFrame,
    quality_summary: Any,
    business_summary: Any,
    kpi_summary: Any,
) -> dict[str, str]:
    date_range = _detect_date_range(current_df)
    return {
        "{{report_title}}": "DataInsight Agent 数据分析 Dashboard",
        "{{date_range}}": date_range,
        "{{row_count}}": f"{len(current_df):,}",
        "{{column_count}}": f"{len(current_df.columns):,}",
        "{{duplicate_count}}": f"{int(current_df.duplicated().sum()):,}",
        "{{missing_count}}": f"{int(current_df.isna().sum().sum()):,}",
        "{{kpi_summary}}": _content_text(kpi_summary, "暂无核心 KPI 结果"),
        "{{business_summary}}": _content_text(
            business_summary,
            "暂无业务分析结果",
        ),
    }


def _replace_dashboard_placeholders(
    worksheet,
    placeholder_values: dict[str, str],
) -> set[str]:
    replaced = set()
    for row in worksheet.iter_rows():
        for cell in row:
            if not isinstance(cell.value, str):
                continue
            updated = cell.value
            for placeholder, replacement in placeholder_values.items():
                if placeholder in updated:
                    updated = updated.replace(placeholder, replacement)
                    replaced.add(placeholder)
            if updated != cell.value:
                cell.value = updated
    return replaced


def _write_dashboard_fallback_summary(
    worksheet,
    placeholder_values: dict[str, str],
    replaced_placeholders: set[str],
) -> None:
    if "{{report_title}}" not in replaced_placeholders:
        worksheet["B2"] = placeholder_values["{{report_title}}"]
    worksheet["B4"] = (
        f"数据更新时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}  |  "
        f"数据范围：{placeholder_values['{{date_range}}']}  |  "
        f"行数：{placeholder_values['{{row_count}}']}  |  "
        f"列数：{placeholder_values['{{column_count}}']}  |  "
        f"重复行：{placeholder_values['{{duplicate_count}}']}  |  "
        f"缺失值：{placeholder_values['{{missing_count}}']}"
    )
    worksheet["B4"].font = Font(size=10, color="666666")
    worksheet["B4"].alignment = Alignment(horizontal="left")

    if "{{row_count}}" not in replaced_placeholders:
        worksheet["C8"] = "数据行数"
        worksheet["C9"] = int(placeholder_values["{{row_count}}"].replace(",", ""))
    if "{{column_count}}" not in replaced_placeholders:
        worksheet["E6"] = "数据列数"
        worksheet["E7"] = int(placeholder_values["{{column_count}}"].replace(",", ""))
    if "{{missing_count}}" not in replaced_placeholders:
        worksheet["H8"] = "缺失值总数"
        worksheet["H9"] = int(placeholder_values["{{missing_count}}"].replace(",", ""))
    if "{{duplicate_count}}" not in replaced_placeholders:
        worksheet["C10"] = f"重复行：{placeholder_values['{{duplicate_count}}']}"


def _detect_date_range(current_df: pd.DataFrame) -> str:
    for column in current_df.columns:
        series = current_df[column]
        if not pd.api.types.is_datetime64_any_dtype(series):
            continue
        valid = series.dropna()
        if valid.empty:
            continue
        return f"{valid.min():%Y-%m-%d} 至 {valid.max():%Y-%m-%d}"
    return "未识别到日期范围"


def _content_text(content: Any, fallback: str) -> str:
    if content is None:
        return fallback
    if isinstance(content, str):
        return content[:2000]
    return json.dumps(content, ensure_ascii=False, default=str)[:2000]


def _build_exploration_summary(current_df: pd.DataFrame) -> pd.DataFrame:
    if current_df.empty:
        return pd.DataFrame({"说明": ["当前数据为空，暂无探索性分析摘要。"]})
    try:
        summary = current_df.describe(include="all", datetime_is_numeric=True).T
    except TypeError:
        summary = current_df.describe(include="all").T
    summary = summary.reset_index().rename(columns={"index": "字段名"})
    return summary


def _append_content_sheet(
    workbook,
    sheet_name: str,
    content: Any,
    fallback: str,
) -> None:
    frame = _as_dataframe(content)
    if frame is None or frame.empty:
        frame = pd.DataFrame({"说明": [fallback]})
    _append_dataframe_sheet(workbook, sheet_name, frame)


def _append_dataframe_sheet(workbook, sheet_name: str, dataframe: pd.DataFrame) -> None:
    worksheet = workbook.create_sheet(sheet_name[:31])
    frame = _excel_safe_dataframe(dataframe).copy()
    if frame.empty and len(frame.columns) == 0:
        frame = pd.DataFrame({"说明": ["暂无数据"]})
    for row in dataframe_to_rows(frame, index=False, header=True):
        worksheet.append(list(row))

    header_fill = PatternFill("solid", fgColor="1F4E78")
    for cell in worksheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
    worksheet.freeze_panes = "A2"
    worksheet.auto_filter.ref = worksheet.dimensions

    for column_index, column_cells in enumerate(
        worksheet.iter_cols(min_row=1, max_row=min(worksheet.max_row, 200)),
        start=1,
    ):
        longest = max((len(str(cell.value or "")) for cell in column_cells), default=8)
        worksheet.column_dimensions[get_column_letter(column_index)].width = min(
            max(longest + 2, 12),
            36,
        )


def _as_dataframe(content: Any) -> pd.DataFrame | None:
    if content is None:
        return None
    if isinstance(content, pd.DataFrame):
        return content.copy()
    if isinstance(content, dict):
        rows = []
        for key, value in content.items():
            if isinstance(value, (list, tuple, set, dict)):
                value = json.dumps(value, ensure_ascii=False, default=str)
            rows.append({"指标": key, "结果": value})
        return pd.DataFrame(rows)
    if isinstance(content, list):
        return pd.DataFrame(content)
    return pd.DataFrame({"结果": [content]})


def _excel_safe_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    result = df.copy()
    for column in result.columns:
        if isinstance(result[column].dtype, pd.DatetimeTZDtype):
            result[column] = result[column].dt.tz_localize(None)
        elif result[column].dtype == "object":
            result[column] = result[column].map(
                lambda value: json.dumps(value, ensure_ascii=False, default=str)
                if isinstance(value, (dict, list, tuple, set))
                else value
            )
    return result


def _add_content(document: Document, content: Any) -> None:
    frame = _as_dataframe(content)
    if frame is None or frame.empty:
        document.add_paragraph("暂无可用结果。")
        return
    _add_docx_table(document, frame)


def _add_docx_table(document: Document, dataframe: pd.DataFrame, max_rows: int = 30) -> None:
    frame = dataframe.head(max_rows).fillna("").astype(str)
    table = document.add_table(rows=1, cols=len(frame.columns))
    table.style = "Table Grid"
    for index, column in enumerate(frame.columns):
        table.rows[0].cells[index].text = str(column)
    for row in frame.itertuples(index=False, name=None):
        cells = table.add_row().cells
        for index, value in enumerate(row):
            cells[index].text = str(value)
    if len(dataframe) > max_rows:
        document.add_paragraph(f"表格仅展示前 {max_rows} 行，共 {len(dataframe)} 行。")


def _pptx_dependencies():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("生成 PPT 需要安装 python-pptx。") from exc
    return Presentation, Inches, Pt


def _add_ppt_title_slide(presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_ppt_bullets(presentation, title: str, lines: list[str], Pt) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def _save_presentation(presentation) -> bytes:
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _content_lines(content: Any, limit: int = 10) -> list[str]:
    frame = _as_dataframe(content)
    if frame is None or frame.empty:
        return []
    lines = []
    for row in frame.head(limit).to_dict("records"):
        line = "；".join(f"{key}：{value}" for key, value in row.items())
        lines.append(line if len(line) <= 240 else f"{line[:237]}...")
    return lines


def _text_lines(text: str | None, limit: int = 10) -> list[str]:
    if not text:
        return []
    return [line.strip(" -#*0123456789.、") for line in text.splitlines() if line.strip()][:limit]
